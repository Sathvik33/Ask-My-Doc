import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer, pipeline
import faiss
import numpy as np
import torch
import os


embedder = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')

model_name = "google/flan-t5-large"
tokenizer = AutoTokenizer.from_pretrained(model_name)
model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
qa_pipeline = pipeline(
    "text2text-generation",
    model=model,
    tokenizer=tokenizer,
    device=0 if torch.cuda.is_available() else -1
)

def read_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        if page.extract_text():
            text += page.extract_text() + "\n"
    return text

def read_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip() != ""])
    return text

def read_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def create_index(text, chunk_size=500):
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    embeddings = embedder.encode(chunks, convert_to_numpy=True)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(embeddings)
    return index, chunks

def query_index(question, index, chunks, top_k=3):
    q_emb = embedder.encode([question], convert_to_numpy=True)
    D, I = index.search(q_emb, top_k)
    context = " ".join([chunks[i] for i in I[0]])
    prompt = f"""
    Answer the question based only on this context:

    {context}

    Question: {question}

    Format your answer as clear, concise bullet points. 
    If possible, number the points.
    """
    answer = qa_pipeline(prompt, max_length=256)[0]['generated_text']
    return answer

st.title("Q&A Bot")

file = st.file_uploader("Upload a file", type=["pdf", "docx", "pptx"])

if file:
    ext = os.path.splitext(file.name)[1].lower()

    with st.spinner("Reading and indexing file..."):
        if ext == ".pdf":
            text = read_pdf(file)
        elif ext == ".docx":
            text = read_docx(file)
        elif ext == ".pptx":
            text = read_pptx(file)
        else:
            st.error(" Unsupported file format")
            st.stop()

        index, chunks = create_index(text)

    st.success(f"{file.name} processed! Ask me anything about it below.")

    question = st.text_input("Enter your question:")

    if st.button("Get Answer") and question:
        with st.spinner("Thinking..."):
            answer = query_index(question, index, chunks)
        st.subheader("Answer:")
        st.write(answer)

